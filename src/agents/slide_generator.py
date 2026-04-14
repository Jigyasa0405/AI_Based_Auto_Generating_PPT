import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Theme Palettes ────────────────────────────────────────────────────────────
THEMES = {
    "blue": {
        "primary":    RGBColor(0x1A, 0x37, 0x6C),
        "accent":     RGBColor(0x23, 0x8B, 0xE6),
        "light_bg":   RGBColor(0xF0, 0xF6, 0xFF),
        "text_dark":  RGBColor(0x0D, 0x1B, 0x2A),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bullet_dot": RGBColor(0x23, 0x8B, 0xE6),
    },
    "green": {
        "primary":    RGBColor(0x0D, 0x47, 0x2F),
        "accent":     RGBColor(0x2E, 0xCC, 0x71),
        "light_bg":   RGBColor(0xF0, 0xFB, 0xF4),
        "text_dark":  RGBColor(0x0A, 0x2E, 0x1E),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bullet_dot": RGBColor(0x2E, 0xCC, 0x71),
    },
    "purple": {
        "primary":    RGBColor(0x2D, 0x10, 0x6B),
        "accent":     RGBColor(0x9B, 0x59, 0xB6),
        "light_bg":   RGBColor(0xF7, 0xF0, 0xFF),
        "text_dark":  RGBColor(0x1A, 0x08, 0x3D),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bullet_dot": RGBColor(0x9B, 0x59, 0xB6),
    },
    "orange": {
        "primary":    RGBColor(0x7A, 0x2E, 0x00),
        "accent":     RGBColor(0xF3, 0x97, 0x1C),
        "light_bg":   RGBColor(0xFF, 0xF8, 0xF0),
        "text_dark":  RGBColor(0x3D, 0x17, 0x00),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bullet_dot": RGBColor(0xF3, 0x97, 0x1C),
    },
    "red": {
        "primary":    RGBColor(0x6B, 0x0F, 0x1A),
        "accent":     RGBColor(0xE7, 0x3C, 0x3E),
        "light_bg":   RGBColor(0xFF, 0xF0, 0xF0),
        "text_dark":  RGBColor(0x3D, 0x07, 0x0E),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bullet_dot": RGBColor(0xE7, 0x3C, 0x3E),
    },
    "teal": {
        "primary":    RGBColor(0x08, 0x46, 0x52),
        "accent":     RGBColor(0x1A, 0xBC, 0x9C),
        "light_bg":   RGBColor(0xEF, 0xFB, 0xF9),
        "text_dark":  RGBColor(0x05, 0x2B, 0x33),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bullet_dot": RGBColor(0x1A, 0xBC, 0x9C),
    },
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ──────────────────────────────────────────────────────────

def _fill_shape(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    _fill_shape(shape, color)
    _no_line(shape)
    return shape


def _add_textbox(slide, left, top, width, height, text, size, bold, color,
                 align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


# ── SlideGenerator ─────────────────────────────────────────────────────────────

class SlideGenerator:
    def __init__(self, output_dir: str = "output/"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def _get_theme(self, name: str) -> dict:
        return THEMES.get(name, THEMES["blue"])

    # ── Cover ──────────────────────────────────────────────────────────────────
    def _make_cover(self, prs, title: str, subtitle: str, theme: dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W, H = SLIDE_W, SLIDE_H

        _add_rect(slide, 0, 0, W, H, theme["primary"])
        _add_rect(slide, W - Inches(1.5), 0, Inches(1.5), H, theme["accent"])
        _add_rect(slide, 0, H - Inches(0.6), W - Inches(1.5), Inches(0.6), theme["accent"])

        _add_textbox(slide, Inches(0.8), Inches(1.8), Inches(10.5), Inches(2.5),
                     title, 44, True, theme["text_light"])
        _add_textbox(slide, Inches(0.8), Inches(4.4), Inches(10.0), Inches(1.0),
                     subtitle, 20, False, theme["accent"])
        _add_textbox(slide, Inches(0.8), H - Inches(0.9), Inches(6), Inches(0.5),
                     "AI-Generated Presentation", 11, False, RGBColor(0xCC, 0xDD, 0xFF))

    # ── Content slide with image: bullets left, image right ────────────────────
    def _make_content_slide_with_image(self, prs, slide_title: str,
                                       bullet_points: list, theme: dict,
                                       image_path: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W, H = SLIDE_W, SLIDE_H

        # Background
        _add_rect(slide, 0, 0, W, H, theme["light_bg"])

        # Right panel (image area) — 44% width
        img_left   = Inches(7.5)
        img_top    = Inches(1.3)
        img_width  = W - img_left
        img_height = H - Inches(1.3)

        # Dark backdrop for image panel
        _add_rect(slide, img_left, img_top, img_width, img_height, theme["primary"])

        # Try embedding the image
        try:
            slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)
            # Thin accent divider
            _add_rect(slide, img_left - Inches(0.07), img_top,
                      Inches(0.07), img_height, theme["accent"])
        except Exception as e:
            print(f"    ⚠  Image embed failed ({e}), using text-only layout.")
            # Fall back gracefully: replace image panel with accent fill
            _add_rect(slide, img_left, img_top, img_width, img_height, theme["accent"])
            _add_textbox(slide, img_left + Inches(0.3), img_top + Inches(1.5),
                         img_width - Inches(0.6), Inches(1.0),
                         "[ Image unavailable ]", 14, False,
                         theme["text_light"], PP_ALIGN.CENTER, italic=True)

        # Header bar over left panel only
        _add_rect(slide, 0, 0, img_left, Inches(1.3), theme["primary"])
        _add_rect(slide, 0, Inches(1.3), Inches(0.12), H - Inches(1.3), theme["accent"])

        # Slide title
        _add_textbox(slide, Inches(0.4), Inches(0.18), img_left - Inches(0.5), Inches(0.95),
                     slide_title, 22, True, theme["text_light"])

        # Bullet points (left panel)
        top = Inches(1.6)
        spacing = Inches(0.80)
        for i, point in enumerate(bullet_points[:5]):
            y = top + i * spacing
            dot = slide.shapes.add_shape(9,
                Inches(0.3), y + Inches(0.13), Inches(0.16), Inches(0.16))
            _fill_shape(dot, theme["bullet_dot"])
            _no_line(dot)

            tb = slide.shapes.add_textbox(Inches(0.60), y,
                                          img_left - Inches(0.75), Inches(0.68))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = point
            run.font.size = Pt(15)
            run.font.color.rgb = theme["text_dark"]
            run.font.name = "Calibri"

    # ── Content slide — text only (full width) ─────────────────────────────────
    def _make_content_slide(self, prs, slide_title: str,
                            bullet_points: list, theme: dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W, H = SLIDE_W, SLIDE_H

        _add_rect(slide, 0, 0, W, H, theme["light_bg"])
        _add_rect(slide, 0, 0, W, Inches(1.3), theme["primary"])
        _add_rect(slide, 0, Inches(1.3), Inches(0.12), H - Inches(1.3), theme["accent"])

        _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(11.5), Inches(0.9),
                     slide_title, 26, True, theme["text_light"])

        top = Inches(1.65)
        spacing = Inches(0.72)
        for i, point in enumerate(bullet_points[:6]):
            y = top + i * spacing
            dot = slide.shapes.add_shape(9,
                Inches(0.35), y + Inches(0.12), Inches(0.18), Inches(0.18))
            _fill_shape(dot, theme["bullet_dot"])
            _no_line(dot)

            tb = slide.shapes.add_textbox(Inches(0.65), y, Inches(11.8), Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = point
            run.font.size = Pt(17)
            run.font.color.rgb = theme["text_dark"]
            run.font.name = "Calibri"

    # ── Conclusion ─────────────────────────────────────────────────────────────
    def _make_conclusion_slide(self, prs, slide_title: str,
                               bullet_points: list, theme: dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W, H = SLIDE_W, SLIDE_H

        _add_rect(slide, 0, 0, W, H, theme["primary"])
        _add_rect(slide, 0, H - Inches(0.5), W, Inches(0.5), theme["accent"])

        _add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1.0),
                     slide_title, 30, True, theme["text_light"])
        _add_rect(slide, Inches(0.8), Inches(1.55), Inches(3.5), Inches(0.06), theme["accent"])

        top = Inches(1.9)
        spacing = Inches(0.75)
        for i, point in enumerate(bullet_points[:5]):
            y = top + i * spacing
            _add_rect(slide, Inches(0.8), y + Inches(0.1),
                      Inches(0.06), Inches(0.38), theme["accent"])
            tb = slide.shapes.add_textbox(Inches(1.1), y, Inches(10.8), Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = point
            run.font.size = Pt(17)
            run.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
            run.font.name = "Calibri"

        _add_textbox(slide, Inches(0.8), H - Inches(1.3), Inches(6), Inches(0.5),
                     "Thank you", 14, True, theme["accent"])

    # ── Public API ─────────────────────────────────────────────────────────────

    def create_slide_deck(self, structured_data: dict,
                          output_filename: str = "generated_presentation.pptx",
                          image_map: dict = None) -> str:
        """
        Build a styled PowerPoint.

        Args:
            structured_data: dict with keys: presentation_title, subtitle,
                             theme_color, slides (list of slide dicts).
            output_filename:  .pptx filename saved inside self.output_dir
            image_map:        optional {slide_index: image_file_path}.
                              Slides with an entry use the split image layout.
        """
        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H

        image_map = image_map or {}
        theme    = self._get_theme(structured_data.get("theme_color", "blue"))
        title    = structured_data.get("presentation_title", "Presentation")
        subtitle = structured_data.get("subtitle", "")
        slides   = structured_data.get("slides", [])

        # 1. Cover slide
        self._make_cover(prs, title, subtitle, theme)

        # 2. Content slides
        for i, slide_data in enumerate(slides):
            slide_type  = slide_data.get("slide_type", "content")
            slide_title = slide_data.get("slide_title", "")
            bullets     = slide_data.get("bullet_points", [])
            note        = slide_data.get("speaker_note", "")

            if slide_type == "conclusion":
                self._make_conclusion_slide(prs, slide_title, bullets, theme)
            elif i in image_map and image_map[i]:
                self._make_content_slide_with_image(
                    prs, slide_title, bullets, theme, image_map[i])
            else:
                self._make_content_slide(prs, slide_title, bullets, theme)

            # Speaker notes
            if note:
                try:
                    prs.slides[-1].notes_slide.notes_text_frame.text = note
                except Exception:
                    pass

        pptx_path = os.path.join(self.output_dir, output_filename)
        prs.save(pptx_path)
        print(f"  ✓ Presentation saved: {pptx_path}")
        return pptx_path


# ── Standalone test ────────────────────────────────────────────────────────────
if __name__ == "__main__":
    from src.agents.image_generator import ImageGenerator

    sample_data = {
        "presentation_title": "AI-Based Slide Generation",
        "subtitle": "Automating beautiful presentations",
        "theme_color": "teal",
        "slides": [
            {
                "slide_title": "What is This Project?",
                "slide_type": "intro",
                "bullet_points": [
                    "Automatically generates PowerPoint slides from documents",
                    "Supports PDF, DOCX, TXT, and CSV input files",
                    "Uses LLaMA AI model via Groq for summarization",
                ],
                "speaker_note": "Introduce the project.",
            },
            {
                "slide_title": "How It Works",
                "slide_type": "content",
                "bullet_points": [
                    "Upload your document via the web interface",
                    "AI extracts and processes the text content",
                    "LLaMA summarizes content into slide-ready bullet points",
                    "Styled PowerPoint is generated and ready to download",
                ],
                "speaker_note": "Walk through the pipeline.",
            },
            {
                "slide_title": "Key Benefits",
                "slide_type": "content",
                "bullet_points": [
                    "Saves hours of manual slide creation time",
                    "Consistent professional visual styling",
                    "100% free — no paid API keys required",
                ],
                "speaker_note": "Emphasize value.",
            },
            {
                "slide_title": "Conclusion",
                "slide_type": "conclusion",
                "bullet_points": [
                    "AI-powered slide generation is practical and effective",
                    "Future: chart support and multi-language documents",
                ],
                "speaker_note": "Wrap up.",
            },
        ],
    }

    img_gen   = ImageGenerator()
    image_map = img_gen.get_images_for_presentation(sample_data["slides"])
    print(f"Images fetched for slide indices: {list(image_map.keys())}\n")

    gen  = SlideGenerator(output_dir="output/")
    path = gen.create_slide_deck(sample_data, image_map=image_map)
    print(f"Done! Open: {path}")